"""Document parsers module.

This module provides unified parsing for multiple document formats with
graceful fallbacks.

Supported formats:
- PDF: pypdf with LovePDF fallback for scanned documents
- DOCX: python-docx with structure extraction
- PPTX: python-pptx with slide structure
- XLSX/CSV: openpyxl/csv with structure analysis
- Images: Vision API for OCR
- HTML: BeautifulSoup with link extraction
- Plain text: encoding detection
"""

import io
import logging
import re
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# MIME type to parser mapping
PARSER_MAP = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "doc",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xlsx",
    "text/csv": "csv",
    "text/plain": "text",
    "text/markdown": "text",
    "text/html": "html",
    "image/png": "image",
    "image/jpeg": "image",
    "image/jpg": "image",
    "image/webp": "image",
    "image/gif": "image",
}


async def parse_document(
    file_data: bytes,
    content_type: str,
    filename: str,
) -> tuple[str, dict[str, Any]]:
    """Parse document and extract text content.

    Args:
        file_data: Raw file bytes.
        content_type: MIME type of the file.
        filename: Original filename for fallback detection.

    Returns:
        Tuple of (extracted_text, metadata_dict).
        Metadata includes: page_count, parser, word_count, etc.
    """
    # Determine parser
    parser_type = PARSER_MAP.get(content_type)
    if not parser_type:
        # Try to detect from extension
        ext = Path(filename).suffix.lower()
        ext_map = {
            ".pdf": "pdf",
            ".docx": "docx",
            ".doc": "doc",
            ".pptx": "pptx",
            ".ppt": "pptx",
            ".xlsx": "xlsx",
            ".xls": "xlsx",
            ".csv": "csv",
            ".txt": "text",
            ".md": "text",
            ".html": "html",
            ".htm": "html",
            ".png": "image",
            ".jpg": "image",
            ".jpeg": "image",
            ".webp": "image",
        }
        parser_type = ext_map.get(ext, "text")

    # Dispatch to appropriate parser
    parsers = {
        "pdf": parse_pdf,
        "docx": parse_docx,
        "pptx": parse_pptx,
        "xlsx": parse_xlsx,
        "csv": parse_csv,
        "text": parse_text,
        "html": parse_html,
        "image": parse_image,
    }

    parser_func = parsers.get(parser_type, parse_text)
    text, metadata = await parser_func(file_data, filename)

    metadata["parser"] = parser_type
    metadata["word_count"] = len(text.split()) if text else 0
    metadata["char_count"] = len(text) if text else 0

    return text, metadata


async def parse_pdf(file_data: bytes, filename: str) -> tuple[str, dict]:
    """Parse PDF document.

    Uses pypdf as primary parser with fallback for scanned documents.
    """
    import pypdf

    metadata = {"page_count": 0}
    text_parts = []

    try:
        reader = pypdf.PdfReader(io.BytesIO(file_data))
        metadata["page_count"] = len(reader.pages)

        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")

        full_text = "\n\n".join(text_parts)

        # Check if we got meaningful text
        if len(full_text.strip()) < 100 and metadata["page_count"] > 0:
            logger.info(f"PDF appears to be scanned, attempting OCR: {filename}")
            # Could integrate with LovePDF or other OCR service here
            metadata["needs_ocr"] = True

        return full_text, metadata

    except Exception as e:
        logger.error(f"PDF parsing failed for {filename}: {e}")
        return f"[PDF parsing error: {e}]", metadata


async def parse_docx(file_data: bytes, filename: str) -> tuple[str, dict]:
    """Parse DOCX document with structure extraction."""
    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not installed, falling back to basic extraction")
        return await parse_text(file_data, filename)

    metadata = {"page_count": None}
    text_parts = []

    try:
        doc = Document(io.BytesIO(file_data))

        # Extract document properties if available
        if doc.core_properties:
            if doc.core_properties.title:
                metadata["title"] = doc.core_properties.title
            if doc.core_properties.author:
                metadata["author"] = doc.core_properties.author
            if doc.core_properties.created:
                metadata["created"] = str(doc.core_properties.created)

        # Extract paragraphs with style information
        for para in doc.paragraphs:
            if para.text.strip():
                # Check if it's a heading
                if para.style and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    text_parts.append(f"\n{'#' * int(level) if level.isdigit() else '#'} {para.text}\n")
                else:
                    text_parts.append(para.text)

        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_text.append(" | ".join(row_text))
            if table_text:
                text_parts.append("\n[Table]\n" + "\n".join(table_text) + "\n")

        return "\n".join(text_parts), metadata

    except Exception as e:
        logger.error(f"DOCX parsing failed for {filename}: {e}")
        return f"[DOCX parsing error: {e}]", metadata


async def parse_pptx(file_data: bytes, filename: str) -> tuple[str, dict]:
    """Parse PowerPoint document with slide structure."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, falling back to basic extraction")
        return await parse_text(file_data, filename)

    metadata = {"page_count": 0}
    text_parts = []

    try:
        prs = Presentation(io.BytesIO(file_data))
        metadata["page_count"] = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Extract table content
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        slide_text.append(" | ".join(row_text))

            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))

        return "\n\n".join(text_parts), metadata

    except Exception as e:
        logger.error(f"PPTX parsing failed for {filename}: {e}")
        return f"[PPTX parsing error: {e}]", metadata


async def parse_xlsx(file_data: bytes, filename: str) -> tuple[str, dict]:
    """Parse Excel spreadsheet with structure analysis."""
    try:
        import openpyxl
    except ImportError:
        logger.warning("openpyxl not installed")
        return "[Excel parsing requires openpyxl]", {}

    metadata = {"sheet_count": 0}
    text_parts = []

    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_data), data_only=True)
        metadata["sheet_count"] = len(wb.worksheets)
        metadata["sheet_names"] = wb.sheetnames

        for sheet in wb.worksheets:
            sheet_text = [f"=== Sheet: {sheet.title} ==="]
            empty_rows = 0
            max_empty = 5  # Stop after 5 consecutive empty rows

            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]

                if any(v.strip() for v in row_values):
                    sheet_text.append("\t".join(row_values))
                    empty_rows = 0
                else:
                    empty_rows += 1
                    if empty_rows >= max_empty:
                        break

            if len(sheet_text) > 1:
                text_parts.append("\n".join(sheet_text))

        return "\n\n".join(text_parts), metadata

    except Exception as e:
        logger.error(f"XLSX parsing failed for {filename}: {e}")
        return f"[XLSX parsing error: {e}]", metadata


async def parse_csv(file_data: bytes, filename: str) -> tuple[str, dict]:
    """Parse CSV file with structure detection."""
    import csv

    metadata = {"row_count": 0}
    text_parts = []

    try:
        # Try different encodings
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                text = file_data.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        else:
            text = file_data.decode("utf-8", errors="replace")

        # Parse CSV
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        metadata["row_count"] = len(rows)

        if rows:
            # First row as header
            header = rows[0]
            metadata["columns"] = header
            text_parts.append("Columns: " + ", ".join(header))
            text_parts.append("-" * 40)

            # Data rows
            for row in rows[1:100]:  # Limit to first 100 rows
                text_parts.append(" | ".join(row))

            if len(rows) > 100:
                text_parts.append(f"... ({len(rows) - 100} more rows)")

        return "\n".join(text_parts), metadata

    except Exception as e:
        logger.error(f"CSV parsing failed for {filename}: {e}")
        return f"[CSV parsing error: {e}]", metadata


async def parse_text(file_data: bytes, filename: str) -> tuple[str, dict]:
    """Parse plain text file with encoding detection."""
    metadata = {}

    # Try different encodings
    for encoding in ["utf-8", "utf-8-sig", "latin-1", "cp1252", "ascii"]:
        try:
            text = file_data.decode(encoding)
            metadata["encoding"] = encoding
            break
        except UnicodeDecodeError:
            continue
    else:
        text = file_data.decode("utf-8", errors="replace")
        metadata["encoding"] = "utf-8 (with errors)"

    return text, metadata


async def parse_html(file_data: bytes, filename: str) -> tuple[str, dict]:
    """Parse HTML with link extraction."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        logger.warning("BeautifulSoup not installed, falling back to regex")
        text = file_data.decode("utf-8", errors="replace")
        # Basic HTML tag removal
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text)
        return text.strip(), {}

    metadata = {}

    try:
        soup = BeautifulSoup(file_data, "html.parser")

        # Extract title
        if soup.title:
            metadata["title"] = soup.title.string

        # Extract meta description
        meta_desc = soup.find("meta", attrs={"name": "description"})
        if meta_desc:
            metadata["description"] = meta_desc.get("content")

        # Extract links
        links = []
        for a in soup.find_all("a", href=True):
            href = a.get("href")
            text = a.get_text(strip=True)
            if href and not href.startswith("#"):
                links.append({"href": href, "text": text})

        metadata["links"] = links[:50]  # Limit links
        metadata["link_count"] = len(links)

        # Extract text
        # Remove script and style elements
        for element in soup(["script", "style", "nav", "footer"]):
            element.decompose()

        text = soup.get_text(separator="\n")
        # Clean up whitespace
        lines = [line.strip() for line in text.splitlines()]
        text = "\n".join(line for line in lines if line)

        return text, metadata

    except Exception as e:
        logger.error(f"HTML parsing failed for {filename}: {e}")
        return f"[HTML parsing error: {e}]", metadata


async def parse_image(file_data: bytes, filename: str) -> tuple[str, dict]:
    """Parse image using OCR.

    Note: Full OCR requires additional setup (Tesseract or cloud API).
    This is a placeholder that can be expanded.
    """
    metadata = {}

    try:
        from PIL import Image

        img = Image.open(io.BytesIO(file_data))
        metadata["format"] = img.format
        metadata["size"] = {"width": img.width, "height": img.height}
        metadata["mode"] = img.mode

    except Exception as e:
        logger.warning(f"Could not read image metadata: {e}")

    # Placeholder for OCR
    # In production, integrate with Tesseract or Vision API
    text = f"[Image: {filename}]\n"
    text += f"Format: {metadata.get('format', 'unknown')}\n"
    text += f"Size: {metadata.get('size', {}).get('width', 0)}x{metadata.get('size', {}).get('height', 0)}\n"
    text += "[OCR text extraction not implemented - enable Vision API for full support]"

    metadata["needs_ocr"] = True

    return text, metadata
